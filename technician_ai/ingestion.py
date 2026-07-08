from __future__ import annotations

import base64
import logging
import os
import sys
from pathlib import Path

from docx import Document
from dotenv import load_dotenv
from pptx import Presentation
from pypdf import PdfReader

load_dotenv(override=True)

from . import database as db
from . import retrieval as rag
from . import tagging as tagger

SUPPORTED_EXTS = {".pdf", ".pptx", ".docx", ".xlsx", ".xls"}
log = logging.getLogger(__name__)

USE_VISION_INGEST = os.environ.get("USE_VISION_INGEST", "false").lower() not in ("false", "0", "no")
# Force vision on every page regardless of text quality (for circuit diagrams).
VISION_ALL_PAGES = os.environ.get("VISION_ALL_PAGES", "false").lower() not in ("false", "0", "no")
# Only apply vision to pages within this range, e.g. "1-30". Empty = all pages.
_vision_range_raw = os.environ.get("VISION_PAGE_RANGE", "")
if _vision_range_raw and "-" in _vision_range_raw:
    _lo, _hi = _vision_range_raw.split("-", 1)
    VISION_PAGE_RANGE: tuple[int, int] | None = (int(_lo), int(_hi))
else:
    VISION_PAGE_RANGE = None
# Minimum ratio of "meaningful" words (3+ chars) to total words.
VISION_QUALITY_THRESHOLD = float(os.environ.get("VISION_QUALITY_THRESHOLD", "0.35"))

VISION_PROMPT = (
    "This is page {page_num} of a technical document (machine manual or circuit diagram). "
    "Extract ALL visible information: text labels, numbers, voltage values, current values, "
    "component names, reference numbers (e.g. C102L, C103L), table contents, section headings, "
    "terminal identifiers, power supply specifications, and any other technical data. "
    "Format as structured plain text. Preserve numeric values and units exactly as shown. "
    "If this is a circuit/electrical diagram, pay special attention to power supply voltages, "
    "component labels, and any overview or summary tables."
)


def _fallback_tags(source_label: str) -> dict:
    return {
        "topic_path": [source_label],
        "entry_type": "reference",
        "title": source_label[:120] or "untitled",
    }


def _render_page_as_png(pdf_path: Path, page_num: int) -> bytes:
    """Render a single PDF page to PNG bytes using PyMuPDF (2x scale for quality)."""
    import fitz  # pymupdf
    doc = fitz.open(str(pdf_path))
    page = doc[page_num - 1]
    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
    return pix.tobytes("png")


def _vision_describe_page(pdf_path: Path, page_num: int) -> str:
    """Send a rendered PDF page to the configured LLM vision API and return extracted text."""
    img_bytes = _render_page_as_png(pdf_path, page_num)
    img_b64 = base64.b64encode(img_bytes).decode()
    model = os.environ.get("TECHNICIAN_AI_MODEL", "gpt-4o")
    prompt = VISION_PROMPT.format(page_num=page_num)

    provider = os.environ.get("LLM_PROVIDER", "").lower()
    if not provider:
        if os.environ.get("ANTHROPIC_API_KEY"):
            provider = "anthropic"
        elif os.environ.get("GOOGLE_API_KEY"):
            provider = "google"
        else:
            provider = "openai"

    if provider == "anthropic":
        import anthropic
        client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))
        response = client.messages.create(
            model=model,
            max_tokens=2048,
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                {"type": "text", "text": prompt},
            ]}],
        )
        return response.content[0].text.strip()
    elif provider == "google":
        from google import genai
        from google.genai import types
        client = genai.Client(api_key=os.environ.get("GOOGLE_API_KEY"))
        response = client.models.generate_content(
            model=model,
            config=types.GenerateContentConfig(max_output_tokens=2048),
            contents=[
                types.Part.from_bytes(data=img_bytes, mime_type="image/png"),
                prompt,
            ],
        )
        return response.text.strip()
    else:
        import openai
        client = openai.OpenAI(
            api_key=os.environ.get("OPENAI_API_KEY"),
            base_url=os.environ.get("LLM_BASE_URL") or None,
        )
        response = client.chat.completions.create(
            model=model,
            max_tokens=2048,
            messages=[{"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
                {"type": "text", "text": prompt},
            ]}],
        )
        return response.choices[0].message.content.strip()


def _text_quality(text: str) -> float:
    """Fraction of words that are 3+ characters (meaningful content vs grid noise)."""
    import re
    words = re.findall(r"[A-Za-z0-9一-鿿]{1,}", text)
    if not words:
        return 0.0
    meaningful = sum(1 for w in words if len(w) >= 3)
    return meaningful / len(words)


def _extract_pdf_pages(pdf_path: Path) -> list[tuple[int, str]]:
    reader = PdfReader(str(pdf_path))
    pages: list[tuple[int, str]] = []
    for page_num, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        needs_vision = False
        if USE_VISION_INGEST:
            in_range = (VISION_PAGE_RANGE is None or VISION_PAGE_RANGE[0] <= page_num <= VISION_PAGE_RANGE[1])
            if in_range and (VISION_ALL_PAGES or _text_quality(text) < VISION_QUALITY_THRESHOLD):
                needs_vision = True
                print(f"    p.{page_num}: vision")
        if needs_vision:
            try:
                vision_text = _vision_describe_page(pdf_path, page_num)
                if vision_text:
                    combined = (text + "\n\n[Vision]\n" + vision_text).strip() if text else vision_text
                    pages.append((page_num, combined))
                    continue
            except Exception as e:
                print(f"    p.{page_num}: vision failed ({e}), using text fallback")
        if text:
            pages.append((page_num, text))
    return pages


def _extract_pptx_slides(pptx_path: Path) -> list[tuple[int, str]]:
    prs = Presentation(str(pptx_path))
    slides: list[tuple[int, str]] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes]\n{notes}")
        text = "\n".join(parts).strip()
        if text:
            slides.append((slide_num, text))
    return slides


def _format_table(table) -> str:
    """Render a docx table as a markdown-style text table."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if rows:
        # Insert separator after header row
        widths = [len(c) for c in rows[0].split("|") if c.strip()]
        sep = "| " + " | ".join("-" * max(w, 3) for w in widths) + " |"
        rows.insert(1, sep)
    return "\n".join(rows)


def _extract_docx_sections(docx_path: Path) -> list[tuple[int, str]]:
    """Extract text from a Word document, preserving tables as markdown tables.

    Returns (section_num, text) pairs grouped by heading boundaries so each
    logical section becomes one retrievable chunk.
    """
    doc = Document(str(docx_path))
    sections: list[tuple[int, str]] = []
    current_parts: list[str] = []
    section_num = 1

    def flush():
        nonlocal section_num
        text = "\n\n".join(p for p in current_parts if p).strip()
        if text:
            sections.append((section_num, text))
            section_num += 1
        current_parts.clear()

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if tag == "p":
            from docx.oxml.ns import qn
            from docx.text.paragraph import Paragraph
            para = Paragraph(block, doc)
            style = para.style.name if para.style else ""
            text = para.text.strip()
            if not text:
                continue
            # Flush at each heading to keep sections coherent
            if style.startswith("Heading"):
                flush()
                current_parts.append(f"## {text}")
            else:
                current_parts.append(text)

        elif tag == "tbl":
            from docx.table import Table
            table = Table(block, doc)
            current_parts.append(_format_table(table))

    flush()
    return sections


def _extract_excel_sheets(excel_path: Path) -> list[tuple[int, str]]:
    """Extract each sheet from an Excel file as a markdown table section."""
    import openpyxl
    wb = openpyxl.load_workbook(str(excel_path), data_only=True)
    sections: list[tuple[int, str]] = []
    for sheet_num, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        # Skip entirely empty sheets
        if not any(any(cell is not None for cell in row) for row in rows):
            continue
        # Build markdown table
        lines = [f"## Sheet: {sheet_name}"]
        for i, row in enumerate(rows):
            cells = [str(c) if c is not None else "" for c in row]
            # Drop trailing empty cells
            while cells and cells[-1] == "":
                cells.pop()
            if not cells:
                continue
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("| " + " | ".join("---" for _ in cells) + " |")
        text = "\n".join(lines)
        if text.strip():
            sections.append((sheet_num, text))
    return sections


def ingest_file(
    path: Path,
    organization_id: str | None = None,
    factory_id: str | None = None,
    uploaded_by_user_id: str | None = None,
) -> int:
    if not path.exists():
        raise FileNotFoundError(path)
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"unsupported file type: {ext} (supported: {', '.join(sorted(SUPPORTED_EXTS))})")

    if ext == ".pdf":
        pages = _extract_pdf_pages(path)
        page_label = "page"
    elif ext == ".docx":
        pages = _extract_docx_sections(path)
        page_label = "section"
    elif ext in (".xlsx", ".xls"):
        pages = _extract_excel_sheets(path)
        page_label = "sheet"
    else:
        pages = _extract_pptx_slides(path)
        page_label = "slide"

    title = path.stem
    rows: list[tuple[str, str, list[float] | None, dict]] = []

    page_chunks: list[tuple[int, str]] = []
    for page_num, page_text in pages:
        for chunk in rag.chunk_text(page_text):
            page_chunks.append((page_num, chunk))

    if not page_chunks:
        return 0

    db.init_db()

    use_llm_tagger = os.environ.get("USE_LLM_TAGGER", "true").lower() not in ("false", "0", "no")

    tags_per_chunk: list[dict] = []
    if use_llm_tagger:
        print(f"  tagging {len(page_chunks)} chunks ...")
        existing_topics = db.list_existing_topic_paths(factory_id=factory_id)
        for i, (_, chunk) in enumerate(page_chunks):
            try:
                tags = tagger.tag_content(chunk, source_label=title, existing_topics=existing_topics)
            except Exception as exc:
                log.warning("LLM tagger failed for %s; using fallback tags: %s", title, exc)
                tags = _fallback_tags(title)
            tags_per_chunk.append(tags)
            existing_topics.append(tags["topic_path"])
            if (i + 1) % 5 == 0 or i + 1 == len(page_chunks):
                print(f"    tagged {i + 1}/{len(page_chunks)}")
    else:
        tags_per_chunk = [{"topic_path": [title], "entry_type": "reference", "title": "untitled"}] * len(page_chunks)

    # Then embed (batched if Voyage is enabled).
    if rag.EMBEDDINGS_ENABLED:
        import time
        embeddings: list[list[float] | None] = []
        BATCH = int(os.environ.get("EMBED_BATCH_SIZE", "16"))
        SLEEP = float(os.environ.get("EMBED_BATCH_SLEEP", "0"))
        for i in range(0, len(page_chunks), BATCH):
            batch = [c for _, c in page_chunks[i : i + BATCH]]
            embeddings.extend(rag.embed_texts(batch, input_type="document"))
            if SLEEP > 0 and i + BATCH < len(page_chunks):
                print(f"    embedded {i + BATCH}/{len(page_chunks)}, sleeping {SLEEP}s ...")
                time.sleep(SLEEP)
    else:
        embeddings = [None] * len(page_chunks)

    for (page_num, chunk), tags, embedding in zip(page_chunks, tags_per_chunk, embeddings):
        metadata = {
            "manual_title": title,
            page_label: page_num,
            "source_path": str(path),
            "topic_path": tags["topic_path"],
            "entry_type": tags["entry_type"],
            "title": tags["title"],
        }
        rows.append(("manual_chunk", chunk, embedding, metadata))

    inserted = db.insert_documents_batch(
        rows,
        organization_id=organization_id,
        factory_id=factory_id,
        uploaded_by_user_id=uploaded_by_user_id,
    )
    return len(inserted)


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python ingest.py <path/to/manual.{pdf,pptx}> [more ...]")
        sys.exit(1)

    db.init_db()
    total = 0
    for arg in sys.argv[1:]:
        path = Path(arg)
        print(f"Ingesting {path} ...")
        count = ingest_file(path)
        print(f"  -> {count} chunks")
        total += count
    print(f"Done. {total} chunks total.")


if __name__ == "__main__":
    main()
