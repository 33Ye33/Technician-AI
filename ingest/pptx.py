"""PowerPoint text extraction."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ingest.base import PageResult


def extract(path: Path) -> PageResult:
    prs = Presentation(str(path))
    slides: PageResult = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes]\n{notes}")
        text = "\n".join(parts).strip()
        if text:
            slides.append((slide_num, text))
    return slides
