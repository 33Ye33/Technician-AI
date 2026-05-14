from __future__ import annotations

import logging
import sys
from pathlib import Path

log = logging.getLogger(__name__)

from dotenv import load_dotenv
from pptx import Presentation
from pypdf import PdfReader

load_dotenv()

import db
import rag
import tagger

SUPPORTED_EXTS = {".pdf", ".pptx"}


def _extract_pdf_pages(pdf_path: Path) -> list[tuple[int, str]]:
    reader = PdfReader(str(pdf_path))
    pages: list[tuple[int, str]] = []
    for page_num, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append((page_num, text))
    return pages


def _extract_pptx_slides(pptx_path: Path) -> list[tuple[int, str]]:
    prs = Presentation(str(pptx_path))
    slides: list[tuple[int, str]] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes]\n{notes}")
        text = "\n".join(parts).strip()
        if text:
            slides.append((slide_num, text))
    return slides


def ingest_file(path: Path) -> int:
    if not path.exists():
        raise FileNotFoundError(path)
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"unsupported file type: {ext} (supported: {', '.join(sorted(SUPPORTED_EXTS))})")

    if ext == ".pdf":
        pages = _extract_pdf_pages(path)
        page_label = "page"
    else:
        pages = _extract_pptx_slides(path)
        page_label = "slide"

    title = path.stem
    rows: list[tuple[str, str, list[float] | None, dict]] = []

    page_chunks: list[tuple[int, str]] = []
    for page_num, page_text in pages:
        for chunk in rag.chunk_text(page_text):
            page_chunks.append((page_num, chunk))

    if not page_chunks:
        log.info("[%s] no text extracted, skipping", title)
        return 0

    db.init_db()

    log.info("[%s] tagging %d chunks (batch_size=%d) ...", title, len(page_chunks), tagger.TAG_BATCH_SIZE)
    existing_topics = db.list_existing_topic_paths()
    tags_per_chunk: list[dict] = []
    batch_size = tagger.TAG_BATCH_SIZE
    all_texts = [chunk for _, chunk in page_chunks]
    for i in range(0, len(all_texts), batch_size):
        batch = all_texts[i : i + batch_size]
        batch_tags = tagger.tag_content_batch(batch, source_label=title, existing_topics=existing_topics)
        tags_per_chunk.extend(batch_tags)
        for t in batch_tags:
            existing_topics.append(t["topic_path"])
        log.info("[%s] tagged %d/%d", title, len(tags_per_chunk), len(page_chunks))

    if rag.EMBEDDINGS_ENABLED:
        log.info("[%s] embedding %d chunks ...", title, len(page_chunks))
        embeddings: list[list[float] | None] = []
        BATCH = 64
        for i in range(0, len(page_chunks), BATCH):
            batch = [c for _, c in page_chunks[i : i + BATCH]]
            embeddings.extend(rag.embed_texts(batch, input_type="document"))
        log.info("[%s] embedding complete", title)
    else:
        log.info("[%s] embeddings disabled, skipping", title)
        embeddings = [None] * len(page_chunks)

    for (page_num, chunk), tags, embedding in zip(page_chunks, tags_per_chunk, embeddings):
        metadata = {
            "manual_title": title,
            page_label: page_num,
            "source_path": str(path),
            "topic_path": tags["topic_path"],
            "entry_type": tags["entry_type"],
            "title": tags["title"],
        }
        rows.append(("manual_chunk", chunk, embedding, metadata))

    inserted = db.insert_documents_batch(rows)
    log.info("[%s] inserted %d chunks into db", title, len(inserted))
    return len(inserted)


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python ingest.py <path/to/manual.{pdf,pptx}> [more ...]")
        sys.exit(1)

    db.init_db()
    total = 0
    for arg in sys.argv[1:]:
        path = Path(arg)
        print(f"Ingesting {path} ...")
        count = ingest_file(path)
        print(f"  -> {count} chunks")
        total += count
    print(f"Done. {total} chunks total.")


if __name__ == "__main__":
    main()
